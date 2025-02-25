import asyncio
from googletrans import Translator
import tkinter as tk
from tkinter import filedialog, messagebox
from pptx import Presentation
from serbian_text_converter import SerbianTextConverter

# Google Translate language dictionary
LANGUAGES = {
    "en": "English",
    "sr": "Serbian (Cyrillic)",
    "sr_Latn": "Serbian (Latin)",
    "es": "Spanish",
    "fr": "French",
    "de": "German",
    "it": "Italian",
    "ru": "Russian",
    "pt": "Portuguese",
    "zh-cn": "Chinese (Simplified)",
    "ja": "Japanese",
    "ko": "Korean",
    "ar": "Arabic",
    # Add more languages as necessary
}


class PowerPointTranslationApp:
    def __init__(self, input_path, output_path, target_lang="en"):
        self.input_path = input_path
        self.output_path = output_path
        self.target_language_code = target_lang  # Default to "en" (English)

        # Start the translation process
        self.translate_pptx_file()

    async def translate_text(self, text, translator, target_lang='en'):
        if text and isinstance(text, str):
            try:
                # Translate the text using Google Translate
                translated = await translator.translate(text, dest=target_lang)
                translated_text = translated.text  # Access the .text attribute of the Translated object
                print(f"Original: {text} -> Translated: {translated_text}")

                # If translating to Serbian Latin, convert Cyrillic to Latin
                if target_lang == "sr_Latn":
                    translated_text = SerbianTextConverter.to_latin(translated_text)

                return translated_text
            except Exception as e:
                print(f"Error translating text: {text}, Error: {e}")
                return text
        return text

    async def translate_slide(self, slide, translator, target_lang='en'):
        # Iterate through all shapes in the slide
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:  # If the shape has text
                original_text = shape.text
                # Save the original font properties
                original_font_size = self.get_font_size(shape)

                # Translate the text
                translated_text = await self.translate_text(original_text, translator, target_lang)

                # Reapply the original font size and formatting
                self.set_text_and_formatting(shape, translated_text, original_font_size)

    def get_font_size(self, shape):
        """Extracts and returns the font size of the text in a shape."""
        if hasattr(shape, "text_frame") and shape.text_frame is not None:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    return run.font.size  # Assuming all runs in the shape have the same font size
        return None  # Default if no font size found

    def set_text_and_formatting(self, shape, translated_text, original_font_size):
        """Sets the translated text and applies the original font size."""
        if hasattr(shape, "text_frame") and shape.text_frame is not None:
            # Clear any existing text (needed to reapply formatted text)
            shape.text_frame.clear()

            # Create a new paragraph and set the translated text
            p = shape.text_frame.add_paragraph()
            p.text = translated_text

            # Reapply the original font size if available
            if original_font_size is not None:
                for run in p.runs:
                    run.font.size = original_font_size

    def translate_pptx_file(self):
        if not self.input_path or not self.output_path or not self.target_language_code:
            print("Error: Please provide all required paths and target language.")
            return

        try:
            # Get the current event loop
            loop = asyncio.get_event_loop()

            # Run the async method with the event loop
            loop.run_until_complete(
                self.translate_pptx_file_async(self.input_path, self.output_path, self.target_language_code))
        except Exception as e:
            print(f"Error: An error occurred: {e}")

    async def translate_pptx_file_async(self, input_path, output_path, target_lang):
        try:
            # Load the PowerPoint presentation
            presentation = Presentation(input_path)
            translator = Translator()

            # Process each slide in the presentation
            for slide_number, slide in enumerate(presentation.slides, start=1):
                print(f"Translating slide: {slide_number}")
                await self.translate_slide(slide, translator, target_lang)

            # Save the translated presentation
            presentation.save(output_path)
            print(f"PowerPoint file has been translated and saved as {output_path}.")
        except Exception as e:
            print(f"Error: An error occurred: {e}")


class PowerPointTranslationAppGUI:
    def __init__(self, root):
        self.root = root
        self.root.title("PowerPoint Translator")
        self.root.geometry("700x300")
        self.root.config(bg="#2e2e2e")  # Dark background

        # File path variables
        self.input_path = tk.StringVar()
        self.output_path = tk.StringVar()
        self.output_name = tk.StringVar()

        # Target language variables
        self.target_language_code = tk.StringVar()  # Store language code instead of name
        self.target_language_name = tk.StringVar()  # Store the language name for display
        self.target_language_code.set("en")  # Default to English code (en)
        self.target_language_name.set(LANGUAGES["en"])  # Default to "English" as the name

        # Title label
        self.title_label = tk.Label(self.root, text="PowerPoint Presentation Translator", font=("Arial", 16), fg="white", bg="#2e2e2e")
        self.title_label.pack(pady=10)

        # File selection frame
        self.file_frame = tk.Frame(self.root, bg="#2e2e2e")
        self.file_frame.pack(pady=20)

        self.input_file_label = tk.Label(self.file_frame, text="Select PowerPoint file:", font=("Arial", 12), fg="white", bg="#2e2e2e")
        self.input_file_label.pack(side=tk.LEFT, padx=5)

        self.input_file_button = tk.Button(self.file_frame, text="Browse", command=self.select_file, bg="#4a4a4a", fg="white")
        self.input_file_button.pack(side=tk.LEFT)

        self.input_file_entry = tk.Entry(self.file_frame, textvariable=self.input_path, width=40, font=("Arial", 12))
        self.input_file_entry.pack(side=tk.LEFT, padx=5)

        # Language selection frame
        self.language_frame = tk.Frame(self.root, bg="#2e2e2e")
        self.language_frame.pack(pady=10)

        self.language_label = tk.Label(self.language_frame, text="Select target language:", font=("Arial", 12),
                                       fg="white", bg="#2e2e2e")
        self.language_label.pack(side=tk.LEFT, padx=5)

        # Create a list of language names (for display) and their corresponding language codes (for values)
        language_codes = list(LANGUAGES.keys())  # Language codes for values
        language_names = list(LANGUAGES.values())  # Language names for display

        # Function to update language code based on the selected language name
        def update_language_code(*args):
            selected_language_name = self.target_language_name.get()
            for code, name in LANGUAGES.items():
                if name == selected_language_name:
                    self.target_language_code.set(code)
                    break

        # Trace the target_language_name variable so that update_language_code gets called when the selection changes
        self.target_language_name.trace("w", update_language_code)

        # Create the OptionMenu using the language names (displayed) and corresponding language codes (values)
        self.language_dropdown = tk.OptionMenu(self.language_frame, self.target_language_name, *language_names)
        self.language_dropdown.config(font=("Arial", 12), bg="#4a4a4a", fg="white")
        self.language_dropdown.pack(side=tk.LEFT)

        # Output path and filename
        self.output_frame = tk.Frame(self.root, bg="#2e2e2e")
        self.output_frame.pack(pady=10)

        self.output_label = tk.Label(self.output_frame, text="Select output location and name:", font=("Arial", 12), fg="white", bg="#2e2e2e")
        self.output_label.pack(side=tk.LEFT, padx=5)

        self.output_name_entry = tk.Entry(self.output_frame, textvariable=self.output_name, width=40, font=("Arial", 12))
        self.output_name_entry.pack(side=tk.LEFT, padx=5)

        self.output_path_button = tk.Button(self.output_frame, text="Browse", command=self.select_output_location, bg="#4a4a4a", fg="white")
        self.output_path_button.pack(side=tk.LEFT)

        # Translate button
        self.translate_button = tk.Button(self.root, text="Translate", command=self.translate_pptx, bg="#4a4a4a", fg="white", font=("Arial", 14))
        self.translate_button.pack(pady=20)

    def select_file(self):
        file_path = filedialog.askopenfilename(filetypes=[("PowerPoint Files", "*.pptx")])
        if file_path:
            self.input_path.set(file_path)

            # Automatically update the output path
            output_filename = file_path.split("/")[-1].replace(".pptx", "") + "-translated.pptx"
            output_dir = "/".join(file_path.split("/")[:-1])  # Get the directory of the selected file
            output_path = f"{output_dir}/{output_filename}"
            self.output_path.set(output_path)  # Set the output path

            # Ensure the output text box displays the output path correctly
            self.output_name.set(f"{output_dir}/{output_filename}")  # You can also update the filename part if needed

    def select_output_location(self):
        output_dir = filedialog.askdirectory()
        if output_dir:
            # Default to the same directory as input file if not specified
            output_filename = self.input_path.get().split("/")[-1].replace(".pptx", "") + self.output_name.get()
            self.output_path.set(f"{output_dir}/{output_filename}")

    def translate_pptx(self):
        input_path = self.input_path.get()
        output_path = self.output_path.get()
        target_lang = self.target_language_code.get()

        if not input_path or not output_path:
            messagebox.showerror("Input Error", "Please select input and output files.")
            return

        try:
            # Initialize and start the translation app
            app = PowerPointTranslationApp(input_path, output_path, target_lang)
            messagebox.showinfo("Success", f"Translation completed and saved as {output_path}")
        except Exception as e:
            messagebox.showerror("Translation Error", f"An error occurred during translation: {e}")


if __name__ == "__main__":
    root = tk.Tk()
    app = PowerPointTranslationAppGUI(root)
    root.mainloop()
