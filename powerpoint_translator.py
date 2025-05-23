import threading
from pptx import Presentation
from azure.ai.translation.text import TextTranslationClient
from azure.core.credentials import AzureKeyCredential


class PowerPointTranslationApp:
    def __init__(self, input_path, output_path, target_lang="en", progress_callback=None):
        self.input_path = input_path
        self.output_path = output_path
        self.target_language_code = target_lang  # Default to "en" (English)
        self.progress_callback = progress_callback  # Callback to update progress bar

        # Set up Azure Translator credentials
        self.endpoint = "https://api.cognitive.microsofttranslator.com/"
        key_path = "key.txt"  # The file holding the translator key
        with open(key_path, 'r', encoding='utf-8') as file:
            self.subscription_key = file.read()
        self.region = "westeurope"
        self.client = TextTranslationClient(
            endpoint=self.endpoint,
            credential=AzureKeyCredential(self.subscription_key),
            headers={"Ocp-Apim-Subscription-Region": self.region}
        )

        # Start the translation process
        self.translate_pptx_file()

    def translate_text(self, text, target_lang='en'):
        """Translates the input text."""
        if text and isinstance(text, str):
            try:
                response = self.client.translate(
                    body=[text], to_language=[target_lang]
                )
                translated_text = response[0].translations[0].text
                print(f"Translated text: {translated_text}")

                return translated_text
            except Exception as e:
                print(f"Error translating text: {text}, Error: {e}")
                return text
        return text

    def translate_slide(self, slide, target_lang='en'):
        """Iterate through all shapes in the slide and translate text."""
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame is not None:
                original_text = shape.text_frame.text
                # Save the original font properties
                original_font_size, original_font_color, original_font_bold, original_font_italic, original_font_underline = self.get_text_properties(shape)

                # Translate the text
                translated_text = self.translate_text(original_text, target_lang)

                # Reapply the original formatting
                self.set_text_and_formatting(shape, translated_text, original_font_size, original_font_color,
                                             original_font_bold, original_font_italic, original_font_underline)

    def get_text_properties(self, shape):
        """Extracts and returns the font properties like size, color, bold, italic, and underline."""
        font_size = None
        font_color = None
        font_bold = None
        font_italic = None
        font_underline = None

        if hasattr(shape, "text_frame") and shape.text_frame is not None:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    font_size = run.font.size
                    font_color = run.font.color.rgb if run.font.color else None
                    font_bold = run.font.bold
                    font_italic = run.font.italic
                    font_underline = run.font.underline
                    break  # Assuming all runs have the same properties

        return font_size, font_color, font_bold, font_italic, font_underline

    def set_text_and_formatting(self, shape, translated_text, font_size, font_color, font_bold, font_italic,
                                font_underline):
        """Sets the translated text and applies the original formatting like font size, color, bold, italic, and underline."""
        if hasattr(shape, "text_frame") and shape.text_frame is not None:
            # Clear existing text
            shape.text_frame.clear()

            # Create a new paragraph and set the translated text
            p = shape.text_frame.add_paragraph()
            p.text = translated_text

            # Reapply original font size, color, bold, italic, and underline
            for run in p.runs:
                if font_size is not None:
                    run.font.size = font_size
                if font_color is not None:
                    run.font.color.rgb = font_color
                if font_bold is not None:
                    run.font.bold = font_bold
                if font_italic is not None:
                    run.font.italic = font_italic
                if font_underline is not None:
                    run.font.underline = font_underline

    def process_slide(self, slide, target_lang='en', progress_counter=None, total_slides=None):
        """Translate individual slide using threading for parallel processing."""
        self.translate_slide(slide, target_lang)

        # Update progress after the slide is fully translated
        if progress_counter is not None and total_slides is not None:
            progress_counter[0] += 1  # Increment the counter
            if self.progress_callback:
                progress = int((progress_counter[0] / total_slides) * 100)
                self.progress_callback(progress)

    def translate_pptx_file(self):
        """Translate the entire PowerPoint file."""
        if not self.input_path or not self.output_path or not self.target_language_code:
            print("Error: Please provide all required paths and target language.")
            return

        try:
            # Load the PowerPoint presentation
            presentation = Presentation(self.input_path)

            # Create a list of threads for parallel processing of slides
            threads = []
            total_slides = len(presentation.slides)
            progress_counter = [0]  # Using a list to allow mutable counter in threads

            # Process each slide in the presentation
            for slide_number, slide in enumerate(presentation.slides, start=1):
                print(f"Translating slide: {slide_number}")
                thread = threading.Thread(target=self.process_slide, args=(slide, self.target_language_code, progress_counter, total_slides))
                threads.append(thread)
                thread.start()

            # Wait for all threads to complete
            for thread in threads:
                thread.join()

            # Save the translated presentation
            presentation.save(self.output_path)
            print(f"PowerPoint file has been translated and saved as {self.output_path}.")
        except Exception as e:
            print(f"Error: An error occurred: {e}")
            if self.progress_callback:
                self.progress_callback(0)  # Reset progress in case of error
